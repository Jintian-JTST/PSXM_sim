"""Generate weekly_report_0812.pptx -- dipole shield-radius optimization only."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
FONT = "Arial"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
blank = prs.slide_layouts[6]

MARGIN = Inches(0.7)
TITLE_TOP = Inches(0.45)
TITLE_H = Inches(0.9)
BODY_TOP = Inches(1.5)
BODY_W = SW - 2 * MARGIN


def set_run(r, size, bold=False, color=BLACK, italic=False):
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box


def add_title(slide, text, size=28):
    box = slide.shapes.add_textbox(MARGIN, TITLE_TOP, BODY_W, TITLE_H)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size, bold=True)


def para(tf, text, size=18, bold=False, first=False, space_after=8, color=BLACK, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    set_run(r, size, bold=bold, color=color, italic=italic)


def bullets(tf, items, size=18, level=0, first=False):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = ("      " if level == 1 else "") + ("- " if level == 0 else "* ") + it
        set_run(r, size - (2 if level else 0))


def new_slide(title):
    s = prs.slides.add_slide(blank)
    add_title(s, title)
    return s


def add_block(s, top, kind, payload):
    if kind == 'ul':
        box = add_textbox(s, MARGIN, top, BODY_W, Inches(0.6) * len(payload) + Inches(0.3))
        bullets(box.text_frame, payload, first=True)
        return top + Inches(0.42) * len(payload)
    if kind == 'ul2':
        box = add_textbox(s, MARGIN, top, BODY_W, Inches(0.6) * len(payload) + Inches(0.3))
        bullets(box.text_frame, payload, level=1)
        return top + Inches(0.38) * len(payload)
    if kind == 'img':
        path, caption = payload
        from PIL import Image
        w, h = Image.open(path).size
        pic_h = Inches(3.0)
        pic_w = Emu(pic_h * w / h)
        if pic_w > BODY_W:
            pic_w = BODY_W
            pic_h = Emu(pic_w * h / w)
        x = MARGIN + (BODY_W - pic_w) // 2
        s.shapes.add_picture(path, x, top, width=pic_w, height=pic_h)
        cap = add_textbox(s, x, top + pic_h, pic_w, Inches(0.3))
        cr = cap.text_frame.paragraphs[0].add_run()
        cr.text = caption
        set_run(cr, 12, italic=True, color=GRAY)
        return top + pic_h + Inches(0.5)
    # plain paragraph / bold heading
    box = add_textbox(s, MARGIN, top, BODY_W, Inches(0.6))
    para(box.text_frame, payload, size=18, bold=(kind == 'h'), first=True)
    return top + Inches(0.5)


def build(title, blocks):
    s = new_slide(title)
    top = BODY_TOP
    for kind, payload in blocks:
        top = add_block(s, top, kind, payload)


# ---- 1: title ----
s = new_slide("PSXM - Dipole Shield Radius Optimization")
box = add_textbox(s, MARGIN, Inches(3.0), BODY_W, Inches(2.2))
tf = box.text_frame
para(tf, "Two-Week Progress Report - 2026-07-28 to 2026-08-12", size=22, first=True, space_after=24)
para(tf, "Jintian Wang", size=20, space_after=6)
para(tf, "Group meeting (Node 5) - 2026-08-12", size=18, color=GRAY)

# ---- 2: design question ----
build("Design Question - Choosing the Shield Radius", [
    ('p', "Shield radius is a trade-off:"),
    ('ul', [
        "Closer shield -> better leakage cancellation, but absorbs more of the central field -> higher coil current",
        "Farther shield -> lower current, but more leakage escapes",
    ]),
    ('p', "Goal: find the shield radius that still meets the 1 mT dipole target within the 1000 A coil budget."),
])

# ---- 3: dipole scan result ----
build("Result - Dipole Shield-Radius Scan", [
    ('ul', [
        "Swept shield radius 25-79 mm at a 1 mT dipole target, 1000 A budget",
        "Optimum at the smallest scanned radius: R_sh = 25 mm, required current ~ 222 A",
        "Required current decreases strictly monotonically with radius",
    ]),
    ('img', ("figures/shield_radius_optimization_dipole.png", "Dipole shield-radius scan (required current, shield current, leakage)")),
])

# ---- 4: optimal shield radius ----
build("Optimal Shield Radius", [
    ('ul', [
        "The scan optimum sat at the lower edge -> ask the precise question: with I_coil fixed at 1000 A (hardware cap), what is the minimum radius that still meets the target?",
        "Exact root: R_sh = 22.99 mm",
        "Only 0.49 mm outside the coil ring (22.5 mm) -> mechanically infeasible",
        "Practical minimum: 23.5-24 mm",
    ]),
])

# ---- 5: conclusion ----
build("Conclusion", [
    ('ul', [
        "Best shielding at the smallest feasible radius: R_sh ~ 23.5-24 mm",
        "At R_sh = 25 mm the dipole needs only 222 A - far below the 1000 A cap",
        "The current budget is not the limiting factor; mechanical clearance is",
    ]),
])

out = "weekly_report_0812.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "weekly_report_0812_dipole.pptx"
    prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
